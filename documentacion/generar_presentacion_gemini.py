#!/usr/bin/env python3
"""Presentación breve: Gemini para el analista de sistemas (15 min)."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

HERE = Path(__file__).resolve().parent
OUT = HERE / "Gemini-IA-para-el-Analista-de-Sistemas.pptx"

NAVY = RGBColor(0x1F, 0x4E, 0x79)
NAVY_DARK = RGBColor(0x15, 0x36, 0x54)
BLUE = RGBColor(0x2E, 0x86, 0xAB)
TEAL = RGBColor(0x14, 0x8F, 0x77)
GOLD = RGBColor(0xB7, 0x95, 0x0B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x1C, 0x28, 0x33)
MUTED = RGBColor(0x5D, 0x6D, 0x7E)
CARD = RGBColor(0xEA, 0xF2, 0xF8)
CARD2 = RGBColor(0xE8, 0xF8, 0xF5)
CARD3 = RGBColor(0xFC, 0xF3, 0xCF)
RED_SOFT = RGBColor(0xF5, 0xB7, 0xB1)
LINE = RGBColor(0xD4, 0xE6, 0xF1)


def set_run(run, size=18, bold=False, color=INK, font="Calibri"):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_text(tf, text, size=18, bold=False, color=INK, align=PP_ALIGN.LEFT, space_after=6):
    p = tf.paragraphs[0] if not tf.paragraphs[0].text else tf.add_paragraph()
    if not tf.paragraphs[0].text and len(tf.paragraphs) == 1:
        p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(space_after)
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color)
    return p


def fill_shape(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def rect(slide, l, t, w, h, color):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    fill_shape(sh, color)
    return sh


def round_rect(slide, l, t, w, h, color):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    fill_shape(sh, color)
    return sh


def textbox(slide, l, t, w, h, text, size=18, bold=False, color=INK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color)
    return box


def bullets(slide, l, t, w, h, items, size=18, color=INK, bold_first=False):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(8)
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = "•  " + item
        set_run(run, size=size, bold=(bold_first and i == 0), color=color)
    return box


def footer(slide, page, total=11):
    rect(slide, Inches(0), Inches(7.15), Inches(13.333), Inches(0.35), NAVY)
    textbox(slide, Inches(0.4), Inches(7.16), Inches(10), Inches(0.3),
            "Gemini para el analista de sistemas  ·  Sesión de 15 minutos  ·  Uso interno",
            size=11, color=WHITE)
    textbox(slide, Inches(11.6), Inches(7.16), Inches(1.4), Inches(0.3),
            f"{page}  /  {total}", size=11, color=WHITE, align=PP_ALIGN.RIGHT)


def header_bar(slide, kicker, title):
    rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.15), NAVY)
    rect(slide, Inches(0), Inches(1.15), Inches(13.333), Inches(0.08), TEAL)
    textbox(slide, Inches(0.5), Inches(0.12), Inches(12), Inches(0.3),
            kicker.upper(), size=12, bold=True, color=RGBColor(0xA3, 0xE4, 0xD7))
    textbox(slide, Inches(0.5), Inches(0.42), Inches(12.2), Inches(0.6),
            title, size=26, bold=True, color=WHITE)


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def card(slide, l, t, w, h, title, body, fill=CARD, title_color=NAVY):
    round_rect(slide, l, t, w, h, fill)
    textbox(slide, l + Inches(0.18), t + Inches(0.12), w - Inches(0.3), Inches(0.4),
            title, size=16, bold=True, color=title_color)
    box = slide.shapes.add_textbox(l + Inches(0.18), t + Inches(0.5), w - Inches(0.36), h - Inches(0.62))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = body
    set_run(run, size=14, color=INK)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    total = 11

    # 1 Cover
    s = prs.slides.add_slide(blank)
    rect(s, Inches(0), Inches(0), Inches(13.333), Inches(7.5), NAVY_DARK)
    rect(s, Inches(0), Inches(0), Inches(0.18), Inches(7.5), TEAL)
    textbox(s, Inches(0.7), Inches(1.7), Inches(11.5), Inches(0.4),
            "SESIÓN TÉCNICA  ·  15 MINUTOS", size=14, bold=True, color=RGBColor(0xA3, 0xE4, 0xD7))
    textbox(s, Inches(0.7), Inches(2.15), Inches(12), Inches(1.4),
            "Gemini: una herramienta de trabajo\npara el analista de sistemas",
            size=36, bold=True, color=WHITE)
    textbox(s, Inches(0.7), Inches(4.0), Inches(11.5), Inches(0.8),
            "Qué es, en qué aporta valor y cómo usarla con criterio profesional.\nSin marketing: enfoque operativo, verificable y gobernado.",
            size=18, color=RGBColor(0xD6, 0xEA, 0xF8))
    textbox(s, Inches(0.7), Inches(6.4), Inches(11), Inches(0.4),
            "Área de Tecnología de la Información  ·  Exposición interna",
            size=14, color=RGBColor(0xAE, 0xB6, 0xBF))
    notes(s,
          "Saludo breve (30–40 s). Presentarse como analista de sistemas. Dejar claro el propósito: "
          "no vender Gemini, sino mostrar cuándo conviene usarlo en el trabajo diario y cuándo no. "
          "Anunciar que hay espacio para preguntas al final.")

    # 2 Agenda
    s = prs.slides.add_slide(blank)
    header_bar(s, "Mapa de la sesión", "Qué cubriremos en 15 minutos")
    items = [
        ("01", "3 min", "Qué es Gemini y qué no es"),
        ("02", "4 min", "Valor concreto para el analista de sistemas"),
        ("03", "4 min", "Cómo usarla: método, no magia"),
        ("04", "3 min", "Límites, riesgos y regla de oro"),
        ("05", "1 min", "Cierre y preguntas"),
    ]
    y = 1.5
    for num, tmin, title in items:
        round_rect(s, Inches(0.55), Inches(y), Inches(12.2), Inches(0.85), CARD)
        textbox(s, Inches(0.75), Inches(y + 0.18), Inches(1.1), Inches(0.5), num, size=22, bold=True, color=TEAL)
        textbox(s, Inches(2.0), Inches(y + 0.22), Inches(8.2), Inches(0.45), title, size=20, bold=True, color=NAVY)
        textbox(s, Inches(10.6), Inches(y + 0.24), Inches(1.8), Inches(0.4), tmin, size=16, color=MUTED, align=PP_ALIGN.RIGHT)
        y += 1.0
    footer(s, 2, total)
    notes(s,
          "Leer la agenda en 40 segundos. Enfatizar que la sesión es corta y práctica. "
          "El objetivo del oyente al salir: saber en qué tareas Gemini le ahorra tiempo y qué debe seguir haciendo él.")

    # 3 Qué es
    s = prs.slides.add_slide(blank)
    header_bar(s, "Definición objetiva", "Qué es Gemini — y qué no es")
    card(s, Inches(0.5), Inches(1.5), Inches(6.0), Inches(2.35),
         "Qué es",
         "Familia de modelos de inteligencia artificial de Google. "
         "Procesa texto, código, imágenes y documentos. Se usa en la app Gemini, en Google Workspace (Gmail, Docs, Sheets) y, en entorno empresarial, con controles de datos.",
         CARD)
    card(s, Inches(6.8), Inches(1.5), Inches(6.0), Inches(2.35),
         "Qué no es",
         "No es un experto de la empresa. No conoce solos nuestros procedimientos, inventario ni políticas. "
         "No firma, no aprueba y no reemplaza el criterio del analista ni la evidencia del sistema.",
         RGBColor(0xFD, 0xED, 0xEC), RGBColor(0x92, 0x2B, 0x21))
    card(s, Inches(0.5), Inches(4.05), Inches(4.0), Inches(2.7),
         "Multimodal",
         "Puede leer un procedimiento, un ticket, una captura o un diagrama y devolver un borrador útil para revisar.",
         CARD2, TEAL)
    card(s, Inches(4.7), Inches(4.05), Inches(4.0), Inches(2.7),
         "Contexto largo",
         "Sirve para resumir o contrastar documentos extensos: políticas, actas, bitácoras o especificaciones.",
         CARD, NAVY)
    card(s, Inches(8.9), Inches(4.05), Inches(3.9), Inches(2.7),
         "Asistente, no dueño",
         "El valor está en acelerar el primer borrador. La calidad final la pone el profesional.",
         CARD3, RGBColor(0x7D, 0x66, 0x08))
    footer(s, 3, total)
    notes(s,
          "2–3 min. Evitar jerga. Analogía: Gemini es un analista junior muy rápido al que hay que supervisar. "
          "Mencionar que existe versión de consumo y versión de trabajo (Workspace / Enterprise) con distinta "
          "protección de datos: para temas internos, usar el canal institucional autorizado.")

    # 4 Valor
    s = prs.slides.add_slide(blank)
    header_bar(s, "Por qué usarla", "Dónde aporta valor al analista de sistemas")
    rows = [
        ("Levantamiento", "Convertir entrevistas, correos y notas en casos de uso, reglas de negocio y preguntas pendientes."),
        ("Documentación", "Borradores de procedimientos, actas, checklists y diagramas en lenguaje de la organización."),
        ("Análisis", "Contrastar alternativas, impactos, riesgos y dependencias antes de proponer a Gerencia."),
        ("Soporte y operación", "Resumir tickets, armar hipótesis de causa, guiar checklists de diagnóstico."),
        ("Calidad", "Revisar claridad, omisiones y consistencia entre proceso, inventario y política."),
        ("Transferencia", "Pasar un procedimiento técnico a material de inducción o de exposición, sin perder el fondo."),
    ]
    y = 1.45
    for i, (title, body) in enumerate(rows):
        col = 0 if i % 2 == 0 else 1
        if i % 2 == 0 and i:
            y += 1.75
        x = 0.45 + col * 6.4
        fill = CARD if col == 0 else CARD2
        tc = NAVY if col == 0 else TEAL
        card(s, Inches(x), Inches(y), Inches(6.15), Inches(1.6), title, body, fill, tc)
    footer(s, 4, total)
    notes(s,
          "3 min. Recorrer los seis bloques con un ejemplo de 10 segundos cada uno. "
          "Ejemplo fuerte: documentar un proceso de mantenimiento o de compras a partir de apuntes desordenados. "
          "Dejar el mensaje: Gemini no decide; organiza, redacta y cuestiona para que el analista decida mejor y más rápido.")

    # 5 Casos
    s = prs.slides.add_slide(blank)
    header_bar(s, "Aplicación", "Tres usos que sí valen el tiempo")
    examples = [
        ("1. Del apunte al procedimiento",
         "Entrada: notas de una reunión con TI y Compras.\n"
         "Salida: estructura de procedimiento (objetivo, alcance, pasos, responsables).\n"
         "Valor: de horas de redacción a una revisión de 20–30 minutos."),
        ("2. Del incidente al diagnóstico",
         "Entrada: ticket + síntomas + inventario del equipo.\n"
         "Salida: hipótesis ordenadas, pruebas a ejecutar y criterio de escalamiento.\n"
         "Valor: el técnico no parte de cero; el analista estandariza el método."),
        ("3. Del proceso a la comunicación",
         "Entrada: un procedimiento ya aprobado.\n"
         "Salida: checklist de usuario, correo de convocatoria o diapositiva de 5 puntos.\n"
         "Valor: la misma verdad, en el formato que cada audiencia necesita."),
    ]
    x = 0.4
    for title, body in examples:
        round_rect(s, Inches(x), Inches(1.5), Inches(4.05), Inches(5.2), CARD)
        rect(s, Inches(x), Inches(1.5), Inches(4.05), Inches(0.12), TEAL)
        textbox(s, Inches(x + 0.2), Inches(1.75), Inches(3.65), Inches(1.1),
                title, size=18, bold=True, color=NAVY)
        box = s.shapes.add_textbox(Inches(x + 0.2), Inches(2.9), Inches(3.65), Inches(3.5))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = body
        set_run(run, size=15, color=INK)
        x += 4.25
    footer(s, 5, total)
    notes(s,
          "3 min. Un caso por minuto. En el primero conectar con el trabajo real del área (procedimientos). "
          "En el segundo, insistir: Gemini sugiere pruebas, no cierra el ticket. "
          "En el tercero: comunicación al usuario interno o al guardia/supervisor, en lenguaje simple.")

    # 6 Método
    s = prs.slides.add_slide(blank)
    header_bar(s, "Método de trabajo", "Cómo pedir para obtener trabajo usable")
    steps = [
        ("1. Contexto", "Rol, objetivo y audiencia.\nEj.: «Eres analista de sistemas. Redacta para personal administrativo de Quito»."),
        ("2. Material", "Pegar o adjuntar hechos: inventario, pasos, restricciones, plantilla.\nSin datos, inventará."),
        ("3. Entregable", "Pedir formato exacto: 12 pasos, tabla, correo, riesgos.\nLimitar extensión."),
        ("4. Verificación", "Cruzar con la fuente. Corregir nombres, cargos y reglas.\nUsted publica, no el modelo."),
    ]
    x = 0.4
    for title, body in steps:
        round_rect(s, Inches(x), Inches(1.5), Inches(3.05), Inches(3.55), CARD)
        textbox(s, Inches(x + 0.18), Inches(1.7), Inches(2.7), Inches(0.7), title, size=18, bold=True, color=TEAL)
        box = s.shapes.add_textbox(Inches(x + 0.18), Inches(2.5), Inches(2.7), Inches(2.3))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = body
        set_run(run, size=14, color=INK)
        x += 3.2
    round_rect(s, Inches(0.4), Inches(5.25), Inches(12.5), Inches(1.5), CARD3)
    textbox(s, Inches(0.65), Inches(5.4), Inches(12.1), Inches(0.4),
            "Prompt mínimo que sí funciona", size=16, bold=True, color=RGBColor(0x7D, 0x66, 0x08))
    textbox(s, Inches(0.65), Inches(5.85), Inches(12.1), Inches(0.7),
            "«Con estos apuntes, arma un procedimiento de mantenimiento cada 6 meses, con responsables, excepciones y criterio de derivación a compra. No inventes políticas que no estén en el texto. Señala lo que falta.»",
            size=15, color=INK)
    footer(s, 6, total)
    notes(s,
          "3 min. Insistir en el paso 4: verificación. Leer el prompt de ejemplo. "
          "Decir que pedir ‘hazme un procedimiento bonito’ da texto genérico; pedir restricciones y ‘no inventes’ da un borrador revisable.")

    # 7 Demostración mental / flujo
    s = prs.slides.add_slide(blank)
    header_bar(s, "En la práctica", "Un flujo de 10 minutos en el día a día")
    flow = [
        ("Min 0–2", "Pegar el material\n(ticket, acta, apuntes)"),
        ("Min 2–5", "Pedir el entregable\ncon restricciones"),
        ("Min 5–8", "Contrastar con la\nfuente y la política"),
        ("Min 8–10", "Ajustar, firmar\ny publicar"),
    ]
    x = 0.5
    for i, (t, b) in enumerate(flow):
        round_rect(s, Inches(x), Inches(1.7), Inches(2.7), Inches(2.4), CARD if i % 2 == 0 else CARD2)
        textbox(s, Inches(x + 0.15), Inches(1.9), Inches(2.4), Inches(0.5), t, size=16, bold=True, color=NAVY)
        box = s.shapes.add_textbox(Inches(x + 0.15), Inches(2.5), Inches(2.4), Inches(1.4))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = b
        set_run(run, size=16, color=INK)
        if i < 3:
            textbox(s, Inches(x + 2.55), Inches(2.5), Inches(0.4), Inches(0.5), "→", size=24, bold=True, color=TEAL)
        x += 3.15
    bullets(s, Inches(0.55), Inches(4.4), Inches(12.2), Inches(2.3), [
        "Si el resultado no cita el dato de origen, trátelo como hipótesis, no como hecho.",
        "Si hay cifras, fechas, series o nombres, verifíquelos uno a uno.",
        "Si el tema es confidencial, use solo el entorno autorizado por TI; no pegue contraseñas, datos personales ni secretos de clientes.",
        "Guarde el prompt útil: se vuelve un activo del área, igual que una plantilla.",
    ], size=16)
    footer(s, 7, total)
    notes(s,
          "2 min. Mostrar que 10 minutos bien usados rinden más que 40 minutos de ‘charlar’ con el modelo. "
          "Cerrar con la idea de biblioteca de prompts del área (mantenimiento, compras, tickets, actas).")

    # 8 Riesgos
    s = prs.slides.add_slide(blank)
    header_bar(s, "Gobernanza", "Límites reales — para usarla sin perder control")
    risks = [
        ("Alucinaciones", "Puede afirmar con seguridad algo falso: un paso, una norma o un dato de inventario que no existe."),
        ("Sesgo de fluidez", "Un texto bien escrito no es un texto correcto. La forma no prueba el fondo."),
        ("Datos sensibles", "Cédulas, nómina, credenciales, contratos y datos de clientes no deben ir a canales no autorizados."),
        ("Responsabilidad", "El analista sigue siendo dueño del entregable. Gemini no firma procedimientos ni órdenes de compra."),
    ]
    y = 1.45
    for title, body in risks:
        round_rect(s, Inches(0.5), Inches(y), Inches(12.3), Inches(1.15), RGBColor(0xFD, 0xED, 0xEC))
        textbox(s, Inches(0.75), Inches(y + 0.12), Inches(3.2), Inches(0.85), title, size=18, bold=True, color=RGBColor(0x92, 0x2B, 0x21))
        textbox(s, Inches(4.1), Inches(y + 0.18), Inches(8.4), Inches(0.85), body, size=16, color=INK)
        y += 1.25
    footer(s, 8, total)
    notes(s,
          "2 min. Tono serio, no alarmista. Un ejemplo: si pide un procedimiento y Gemini ‘inventa’ una aprobación de Gerencia que el área no tiene, "
          "eso no se publica. Regla: no hay dato institucional sin fuente.")

    # 9 Regla de oro + checklist
    s = prs.slides.add_slide(blank)
    header_bar(s, "Para llevar", "Regla de oro y checklist antes de publicar")
    round_rect(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.7), TEAL)
    textbox(s, Inches(0.8), Inches(1.7), Inches(11.8), Inches(1.3),
            "Gemini acelera el borrador. El analista garantiza la verdad operativa.\nSi no puede explicarlo o verificarlo, no sale.",
            size=22, bold=True, color=WHITE)
    checks = [
        "¿El texto respeta el procedimiento y los cargos reales del área?",
        "¿Hay algún dato que Gemini pudo haber inventado?",
        "¿Se usó información que no debía salir del canal interno?",
        "¿El entregable queda más claro para el usuario final (admin, guardia, supervisión)?",
        "¿Quedó registrado quién revisó, no solo quién preguntó al modelo?",
    ]
    bullets(s, Inches(0.7), Inches(3.5), Inches(12), Inches(3.2), checks, size=18)
    footer(s, 9, total)
    notes(s,
          "1,5 min. Leer la regla de oro despacio. El checklist puede quedar como práctica del área. "
          "Invitar a usarlo en el próximo procedimiento o informe.")

    # 10 Próximos pasos
    s = prs.slides.add_slide(blank)
    header_bar(s, "Adopción", "Qué hacer esta semana — sin un proyecto grande")
    nexts = [
        ("Hoy", "Elegir una tarea repetitiva: acta, resumen de tickets o primer borrador de un procedimiento."),
        ("Esta semana", "Armar 3 prompts reutilizables del área y probarlos con un caso real, siempre revisando el resultado."),
        ("Este mes", "Acordar con TI qué datos sí / no se pegan en Gemini y en qué cuenta institucional se trabaja."),
        ("No hacer", "No copiar datos personales ni secretos. No publicar salida sin revisión. No sustituir aprobación de Gerencia."),
    ]
    y = 1.5
    colors = [CARD, CARD2, CARD, RGBColor(0xFD, 0xED, 0xEC)]
    tcolors = [NAVY, TEAL, NAVY, RGBColor(0x92, 0x2B, 0x21)]
    for (title, body), fill, tc in zip(nexts, colors, tcolors):
        round_rect(s, Inches(0.5), Inches(y), Inches(12.3), Inches(1.2), fill)
        textbox(s, Inches(0.75), Inches(y + 0.32), Inches(2.4), Inches(0.55), title, size=18, bold=True, color=tc)
        textbox(s, Inches(3.3), Inches(y + 0.28), Inches(9.2), Inches(0.7), body, size=16, color=INK)
        y += 1.3
    footer(s, 10, total)
    notes(s,
          "1 min. Cerrar con acciones pequeñas. El éxito no es ‘usar IA’, es ahorrar tiempo en un entregable real esta semana.")

    # 11 Cierre
    s = prs.slides.add_slide(blank)
    rect(s, Inches(0), Inches(0), Inches(13.333), Inches(7.5), NAVY_DARK)
    rect(s, Inches(0), Inches(0), Inches(0.18), Inches(7.5), TEAL)
    textbox(s, Inches(0.7), Inches(1.6), Inches(12), Inches(0.4),
            "PARA RECORDAR", size=14, bold=True, color=RGBColor(0xA3, 0xE4, 0xD7))
    textbox(s, Inches(0.7), Inches(2.05), Inches(12), Inches(2.2),
            "Gemini no reemplaza al analista.\nLo hace más rápido en el borrador\ny más exigente en la revisión.",
            size=28, bold=True, color=WHITE)
    textbox(s, Inches(0.7), Inches(4.6), Inches(12), Inches(1.0),
            "Tres ideas: contexto + hechos + verificación.\nPreguntas.",
            size=20, color=RGBColor(0xD6, 0xEA, 0xF8))
    textbox(s, Inches(0.7), Inches(6.4), Inches(12), Inches(0.4),
            "Área de Tecnología de la Información",
            size=14, color=RGBColor(0xAE, 0xB6, 0xBF))
    notes(s,
          "Cierre de 45 s y abrir preguntas. Si hay silencio, lanzar: "
          "¿en qué tarea de esta semana les gustaría probar el flujo de 10 minutos?")

    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    print(build())
